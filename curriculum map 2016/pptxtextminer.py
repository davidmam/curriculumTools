from pptx import Presentation
#from docx import Document
import sys
import re
import argparse
# add arguments for module, week etc.
parser = argparse.ArgumentParser(description='Learning event text miner')
parser.add_argument('--module', '-m', help='Module code')
parser.add_argument('--week', '-w', help='Teaching week')
parser.add_argument('--type', '-t', help='Learning event type (Workshop, Lab etc)')
parser.add_argument('--file', '-f', help='Content file (pptx or docx)')
args = parser.parse_args()


anword=re.compile(r'^[0-9a-z-_]*[a-z][0-9a-z-_]*$')

fh=open("/Users/dmamartin/stopwords.txt")
stopwords=[]
for line in fh.readlines():
    stopwords.append(line.strip())

try:
    prs = Presentation(args.file) #take filename
except:
    exit(1)
# text_runs will be populated with a list of strings,
# one for each text run in presentation
text_runs = []

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text_runs.append(run.text)
textbody=" ".join(text_runs).replace("'","").replace('"','')

def get_keywords(text):
    '''simple keywords splitting'''
    kw=filter(lambda x: (anword.match(x) and len(x) >2), filter(lambda y: y not in stopwords, [x.lower() for x in text.split(" ")]))
    return kw

def get_polywords(text):
    '''more complex phrase isolation between stopwords'''
    kw=[x.lower() for x in text.split(" ")]
    for i in range(len(kw)):
        if kw[i] in stopwords:
            kw[i]=","
    kp=[x.strip() for x in " ".join(kw).split(",")]
    
keywords = get_keywords(textbody)

sys.stdout.write(",".join(sorted(keywords)))
