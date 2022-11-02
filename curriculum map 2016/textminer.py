
'''Scripts to read and mine text from teachingmaterials'''

import sys
import os
import re
import argparse
import json
import codecs
sys.path.append('C:\\Python27\\lib\\site-packages')
import PyPDF2
import nltk
import glob
import shutil
from zipfile import ZipFile
from xml.dom.minidom import parse

from bs4 import BeautifulSoup
from pptx import Presentation
from docx import Document
from PyPDF2.utils import  isString, b_, u_, ord_, chr_, str_, formatWarning
from PyPDF2.pdf import ContentStream, TextStringObject

# add arguments for module, week etc.
#parser = argparse.ArgumentParser(description='Learning event text miner')

#parser.add_argument('--dir', '-d', help='Directory to traverse and process all processible files')
#parser.add_argument('--output', '-o', help='Directory in which to store results')
#args = parser.parse_args()


#anword = re.compile(r'^[0-9a-z-_]*[a-z][0-9a-z-_]*$')

#fh = open("stopwords.txt")
#stopwords = []
#for line in fh.readlines():
#    stopwords.append(line.strip())


#textbody = ""

def pdf_textExtract(page, addSpace=False):
    """
    Adapted from the source in PyPDF2 to add spaces between teh text strings. 
    This can break some words but it can fix many more than it breaks, 
    allowing for documents where each word or part  of a word is rendered
    without spaces. addSpace controls this behaviour.
    
    Locate all text drawing commands, in the order they are provided in the
    content stream, and extract the text.  This works well for some PDF
    files, but poorly for others, depending on the generator used.  This will
    be refined in the future.  Do not rely on the order of text coming out of
    this function, as it will change if this function is made more
    sophisticated.
    :return: a unicode string object.
    """
    spacechar = ''
    if addSpace == True:
        spacechar = ' '
    text = u_("")
    content = page["/Contents"].getObject()
    if not isinstance(content, ContentStream):
        content = ContentStream(content, page.pdf)
    # Note: we check all strings are TextStringObjects.  ByteStringObjects
    # are strings where the byte->string encoding was unknown, so adding
    # them to the text here would be gibberish.
    for operands, operator in content.operations:
        if operator == b_("Tj"):
            _text = operands[0]
            if isinstance(_text, TextStringObject):
                text += spacechar + _text
        elif operator == b_("T*"):
            text += "\n"
        elif operator == b_("'"):
            text += "\n"
            _text = operands[0]
            if isinstance(_text, TextStringObject):
                text += spacechar + operands[0]
        elif operator == b_('"'):
            _text = operands[2]
            if isinstance(_text, TextStringObject):
                text += "\n"
                text += _text
        elif operator == b_("TJ"):
            for i in operands[0]:
                if isinstance(i, TextStringObject):
                    text += spacechar + i
            text += "\n"
    return text


def read_pptx(filename):
    '''extract all text strings from a pptx file
    arguments are the filename to read. Returns text in one line'''
    textbody = ''
    try:
        prs = Presentation(filename) #take filename
        # text_runs will be populated with a list of strings,
        # one for each text run in presentation
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    text_runs = []
                    for run in paragraph.runs:
                        text_runs.append(run.text)
                    textbody += ", ".join(text_runs).replace("'", "").replace('"', '')
                    textbody += '\n'
    except:
        pass
    #append notes using Eric Jeng's extract script.
    textbody += read_pptx_notes(filename)
    return textbody

def read_pptx_notes(filename):
    '''Adaptation of Eric Jang's pptx notes reading script.'''
    tmpdir = 'pptxtmp'
    notes=''
    try:
        shutil.rmtree(tmpdir)
    except:
        pass
    try:
        os.mkdir(tmpdir)
    except Exception as e:
        print('Could not create temporary directory for ppt parsing: %s'% e)
        return notes
    ZipFile(filename).extractall(path=tmpdir, pwd=None)
    path = os.path.join(tmpdir, 'ppt','notesSlides')
    xmlfiles = glob.glob(os.path.join(path, '*.xml'))
    for infile in xmlfiles:
        try:
            #parse each XML notes file from the notes folder.
            dom = parse(infile)
            noteslist = dom.getElementsByTagName('a:t')
            if len(noteslist) == 0:
                continue
            for node in noteslist:
                xmlTag = node.toxml()
                notes += xmlTag.replace('<a:t>', ' ').replace('</a:t>', '')
                    #concatenate the xmlData to the tempstring for the particular slideNumber index.
            #os.remove(infile)
        except Exception as e:
            print('Error parsing xml file %s: %s'%(infile,e))
    shutil.rmtree(tmpdir)
    return notes
    
def read_docx(filename):
    '''extract all text strings from a docx file
    arguments are the filename to read. Returns text in one line'''
    textbody = ''
    try:
        document = Document(filename)
        paras = document.paragraphs
        ptlist = []
        for pt in paras:
            ptlist.append(pt.text)
        textbody = "\n".join(ptlist).replace("'", "").replace('"', '')
    except:
        pass
    return textbody

def read_pdf(filename, cutoff=7.0):
    '''extract all text strings from a pdf file
    arguments are the filename to read. Returns text in one line'''
    textbody = ''
    try:
        document = PyPDF2.PdfFileReader(open(filename, 'rb'))
        for pagenum in range(document.numPages):
            page = document.getPage(pagenum)
            rawtext = pdf_textExtract(page, False)
            if len(rawtext) > 0:
                try:
                    tokens = nltk.word_tokenize(rawtext)
                    if len(tokens) > 0:
                        meanlen = len(rawtext)/len(tokens)
                        #print('meanlen is %s'%meanlen)
                        if meanlen > cutoff:
                            rawtext = pdf_textExtract(page, True)
                        #tokens = nltk.word_tokenize(rawtext)
                    #meanlen=len(rawtext)/len(tokens)
                    #print('meanlen is %s'%meanlen)
                except Exception as e:
                    print('Error parsing page %s of %s: %s'%(pagenum, filename, e))
                textbody = textbody + rawtext.replace('"', '').replace("'", "")
    except Exception as e:
        print('Error opening PDF file %s: %s'%(filename, e))
    return textbody

def read_html(filename):
    '''extract all text strings from a html file
    arguments are the filename to read. Returns text in one line'''
    textbody = ''
    try:
        rawtext = str(open(filename, 'rb').read())
        rawtext = re.sub(r'\\x..','', rawtext)
        soup = BeautifulSoup(rawtext, 'html.parser')
        try:        
            rep=soup.script.replace_with(' ')
            while rep:
                rep=soup.script.replace_with(' ')
        except:
            pass
        
        textbody = soup.get_text().replace("'", "").replace('"', '')
    except Exception as e:
        print('Error opening HTML file %s: %s'%(filename, e))
    return textbody

def read_text(filename):
    '''extract text from text file.'''
    return open(filename).read()

def process_file(filename): #TODO
    '''extracts the keywords from the text for a file.'''

    textbody = ''
    if filename.lower()[-4:] == 'pptx':
        textbody = read_pptx(filename)
    elif filename.lower()[-4:] == 'docx':
        textbody = read_docx(filename)
    elif filename.lower()[-3:] == 'pdf':
        textbody = read_pdf(filename)
    elif filename.lower()[-4:] == 'html':
        textbody = read_html(filename)
    elif filename.lower()[-3:] == 'htm':
        textbody = read_html(filename)
    elif filename.lower()[-3:] == 'txt':
        textbody = read_text(filename)
    else:
        print('Incorrect filetype. Cannot parse %s'%filename)
    #keywords = filter(lambda x: (anword.match(x) and len(x) >2),
    #filter(lambda y: y not in stopwords, [x.lower() for x in textbody.split(" ")]))
    return textbody
#    fileroot = filename.split("\\")[-1]
#    eventcode = "_".join([module, week,  fileroot.replace(" ", "_")])
#    sys.stdout.write("%s\t%s\n"%(eventcode, ",".join(sorted(keywords))))

def write_output(outdir, filename, module, week, file, textbody, tfh):
    '''Writes a formatted file to disk containing the annotation for the teaching document
    outdir is the directory in which to save the data
    filename is the filename under which it is to be saved
    module, week and file are the specific document details
    keywords are from the extracted text'''
    fh = open(os.path.join(outdir, filename), 'w')
    fh.write(str(textbody.encode('utf-8')))
    level=module[2]
    semester=module[3]
    tfh.write('\t'.join([filename, level, semester, module, week, file, str(len(textbody))])+'\n')
    fh.close()
    


def read_dir(datasource, datasink):
    '''takes a source directory (datasource) organsied as module/week
    and parses all valid files underneath, storing the results as annotated
    keyword lists in datasink'''
    try:
        tfh = open(os.path.join(datasink, 'index.dat'),'w')
        modules = os.listdir(datasource)
        for m in modules:
            if os.path.isdir(os.path.join(datasource, m)):
                module = m
                mfilecount = 0
                try:
                    weeks = os.listdir(os.path.join(datasource, m))
                    for w in weeks:
                        if os.path.isdir(os.path.join(datasource, m, w)):
                            try:
                                files = os.listdir(os.path.join(datasource, m, w))
                                for f in files:
                                    if f.lower()[-4:] == 'pptx' \
                                            or f.lower()[-4:] == 'docx' \
                                            or f.lower()[-3:] == 'pdf' \
                                            or f.lower()[-4:] == 'html' \
                                            or f.lower()[-3:] == 'htm':
                                        filepath = os.path.join(datasource, m, w, f) 
                                        textbody = process_file(filepath)
                                        mfilecount += 1
                                        fname = '%s_%s.txt'%(module, mfilecount)
                                        write_output(datasink, fname, m, w, f, textbody,tfh)
                            except Exception as e:
                                print('Error reading file %s in week %s of module %s: %s'%(f, w, m, e))
                except Exception as e:
                    print('Error reading events in module %s: %s'%(m, e))
        tfh.close()
    except Exception as e:
        raise Exception('Bad data directory: %s'%e)
        # now annotate and save the data.

#if args.list != None:
#    try:
#        fh = open(args.list)
#        for line in fh.readlines():
#            try:
#                (filename, module, week, ftype) = line.strip().split("\t")
#                process_file(filename, module,week, ftype)
#            except Exception as e:
#                #print "error reading file %s: %s"%(filename, e)
#                continue
#    #process list of files
#    except:
#        print("Error opening %s"%args.list)


#else:
#    docfile = args.file
#    process_file(docfile, args.module, args.week, args.type)
