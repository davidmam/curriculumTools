#TF-IDF calculation. Monolithic python script.

# reads in a data file that is key\tkeywordlist per line.
# keyword list = comma separated terms
import os
import sys
import re
import argparse
import traceback
import math
import nltk
from nltk.corpus import PlaintextCorpusReader, stopwords as STOPWORDS
from nltk.util import ngrams
import pickle
import pyocr
import base64
import pyocr.builders
from PIL import Image
sys.path.append('C:\\Python27\\lib\\site-packages')
import PyPDF2
import glob
import shutil
from zipfile import ZipFile
from xml.dom.minidom import parse

from bs4 import BeautifulSoup
from pptx import Presentation
from docx import Document
from PyPDF2.utils import  isString, b_, u_, ord_, chr_, str_, formatWarning
from PyPDF2.pdf import ContentStream, TextStringObject

# add arguments for module, week etc.
#parser = argparse.ArgumentParser(description='Learning event text miner')

#parser.add_argument('--dir', '-d', help='Directory to traverse and process all processible files')
#parser.add_argument('--output', '-o', help='Directory in which to store results')
#args = parser.parse_args()


#anword = re.compile(r'^[0-9a-z-_]*[a-z][0-9a-z-_]*$')

#fh = open("stopwords.txt")
#stopwords = []
#for line in fh.readlines():
#    stopwords.append(line.strip())


#textbody = ""

def pdf_textExtract(page, addSpace=False):
    """
    Adapted from the source in PyPDF2 to add spaces between teh text strings. 
    This can break some words but it can fix many more than it breaks, 
    allowing for documents where each word or part  of a word is rendered
    without spaces. addSpace controls this behaviour.
    
    Locate all text drawing commands, in the order they are provided in the
    content stream, and extract the text.  This works well for some PDF
    files, but poorly for others, depending on the generator used.  This will
    be refined in the future.  Do not rely on the order of text coming out of
    this function, as it will change if this function is made more
    sophisticated.
    :return: a unicode string object.
    """
    spacechar = ''
    if addSpace == True:
        spacechar = ' '
    text = u_("")
    content = page["/Contents"].getObject()
    if not isinstance(content, ContentStream):
        content = ContentStream(content, page.pdf)
    # Note: we check all strings are TextStringObjects.  ByteStringObjects
    # are strings where the byte->string encoding was unknown, so adding
    # them to the text here would be gibberish.
    for operands, operator in content.operations:
        if operator == b_("Tj"):
            _text = operands[0]
            if isinstance(_text, TextStringObject):
                text += spacechar + _text
        elif operator == b_("T*"):
            text += "\n"
        elif operator == b_("'"):
            text += "\n"
            _text = operands[0]
            if isinstance(_text, TextStringObject):
                text += spacechar + operands[0]
        elif operator == b_('"'):
            _text = operands[2]
            if isinstance(_text, TextStringObject):
                text += "\n"
                text += _text
        elif operator == b_("TJ"):
            for i in operands[0]:
                if isinstance(i, TextStringObject):
                    text += spacechar + i
            text += "\n"
    return text


def read_pptx(filename, images=False):
    '''extract all text strings from a pptx file
    arguments are the filename to read. Returns text in one line'''
    textbody = ''
    try:
        prs = Presentation(filename) #take filename
        # text_runs will be populated with a list of strings,
        # one for each text run in presentation
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    text_runs = []
                    for run in paragraph.runs:
                        try:
                            text_runs.append(run.text)
                        except:
                            pass
                    textbody += ", ".join(text_runs).replace("'", "").replace('"', '')
                    textbody += '\n'
    except Exception as e:
        print('Error reading powerpoint file %s: %s'%(filename, e))
    #append notes using Eric Jeng's extract script.
    textbody += read_pptx_notes(filename)
    if images:
        textbody += read_pptx_images(filename)
    return textbody

def read_pptx_notes(filename):
    '''Adaptation of Eric Jang's pptx notes reading script.'''
    tmpdir = 'pptxtmp'
    notes=''
    try:
        shutil.rmtree(tmpdir)
    except Exception as e:
        print('Error clearing temp file %s: %s'%(tmpdir, e))
    try:
        os.mkdir(tmpdir)
    except Exception as e:
        print('Could not create temporary directory for ppt parsing: %s'% e)
        return notes
    try:    
        ZipFile(filename).extractall(path=tmpdir, pwd=None)
        path = os.path.join(tmpdir, 'ppt','notesSlides')
        xmlfiles = glob.glob(os.path.join(path, '*.xml'))
        for infile in xmlfiles:
            try:
                #parse each XML notes file from the notes folder.
                dom = parse(infile)
                noteslist = dom.getElementsByTagName('a:t')
                if len(noteslist) == 0:
                    continue
                for node in noteslist:
                    xmlTag = node.toxml()
                    notes += xmlTag.replace('<a:t>', ' ').replace('</a:t>', '')
                        #concatenate the xmlData to the tempstring for the particular slideNumber index.
                #os.remove(infile)
            except Exception as e:
                print('Error parsing xml file %s: %s'%(infile,e))
    except:
        print('Error unpacking zip file')
    try:
        shutil.rmtree(tmpdir)
    except Exception as e:
        print('Error clearing up temp file %s: %s'%(tmpdir, e))
    return notes

def read_pptx_images(filename):
    '''Adaptation of Eric Jang's pptx notes reading script.'''
    imagetypes = ('jpg','jpeg','png','tiff','jp2', 'bmp', 'wmf','gif')
    tmpdir = 'pptxtmp'
    notes=''
    try:
        shutil.rmtree(tmpdir)
    except Exception as e:
        print('Error clearing temp file %s: %s'%(tmpdir, e))
    try:
        os.mkdir(tmpdir)
    except Exception as e:
        print('Could not create temporary directory for ppt parsing: %s'% e)
        return notes
    try:    
        ZipFile(filename).extractall(path=tmpdir, pwd=None)
        path = os.path.join(tmpdir, 'ppt','media')
        imagefiles=[ f for f in glob.glob(os.path.join(path, '*.*')) if f.split('.')[-1].lower() in imagetypes ]
        for i in imagefiles:
            txt = read_image_ocr(i)
            notes += ' '+txt
    except Exception as e:
        print('Error extracting text from image %s in %s: %s'%(i,filename, e))
    try:
        shutil.rmtree(tmpdir)
    except Exception as e:
        print('Error clearing up temp file %s: %s'%(tmpdir, e))
    return notes

def read_docx(filename):
    '''extract all text strings from a docx file
    arguments are the filename to read. Returns text in one line'''
    textbody = ''
    try:
        document = Document(filename)
        paras = document.paragraphs
        ptlist = []
        for pt in paras:
            ptlist.append(pt.text)
        textbody = "\n".join(ptlist).replace("'", "").replace('"', '')
    except:
        pass
    return textbody

def read_pdf(filename, cutoff=7.0, do_ocr=False):
    '''extract all text strings from a pdf file
    arguments are the filename to read. Returns text in one line'''
    textbody = ''
    tmpdir = 'tmppdf'
    try:
        document = PyPDF2.PdfFileReader(open(filename, 'rb'))
        for pagenum in range(document.numPages):
            page = document.getPage(pagenum)
            rawtext = pdf_textExtract(page, False)
            if len(rawtext) > 0:
                try:
                    tokens = nltk.word_tokenize(rawtext)
                    if len(tokens) > 0:
                        meanlen = len(rawtext)/len(tokens)
                        #print('meanlen is %s'%meanlen)
                        if meanlen > cutoff:
                            rawtext = pdf_textExtract(page, True)
                        #tokens = nltk.word_tokenize(rawtext)
                    #meanlen=len(rawtext)/len(tokens)
                    #print('meanlen is %s'%meanlen)
                except Exception as e:
                    print('Error parsing page %s of %s: %s'%(pagenum, filename, e))
                textbody = textbody + rawtext.replace('"', '').replace("'", "")
            elif do_ocr: # scanned page maybe
                try:
                    imagefile=''
                    xObject = page['/Resources']['/XObject'].getObject()
                    shutil.rmtree(tmpdir)
                    os.mkdir(tmpdir)    
                    for obj in xObject:
                        if xObject[obj]['/Subtype'] == '/Image':
                            imagefile=''
                            size = (xObject[obj]['/Width'], xObject[obj]['/Height'])
                            data = xObject[obj].getData()
                            if xObject[obj]['/ColorSpace'] == '/DeviceRGB':
                                mode = "RGB"
                            else:
                                mode = "P"
                            imagefile=''
                            if xObject[obj]['/Filter'] == '/FlateDecode':
                                img = Image.frombytes(mode, size, data)
                                imagefile=os.path.join(tmpdir, obj[1:] + ".png")
                                img.save(imagefile)
                            elif xObject[obj]['/Filter'] == '/DCTDecode':
                                imgagefile=os.path.join(tmpdir,obj[1:] + ".jpg")
                                img.open(imagefile, 'wb')
                                img.write(data)
                                img.close()
                            elif xObject[obj]['/Filter'] == '/JPXDecode':
                                imagefile = os.path.join(tmpdir,obj[1:] + ".jp2")
                                img=open(imagefile, "wb")
                                img.write(data)
                                img.close()
                            ocrtext = read_image_ocr(imagefile)
                            if len(ocrtext) > 10:
                                textbody = textbody + ' ' + ocrtext
                    shutil.rmtree(tmpdir)
                except Exception as e:
                    print('Error reading imagefile %s in %s: %s'%(imagefile, filename, e) )
    except Exception as e:
        print('Error opening PDF file %s: %s'%(filename, e))
    textbody = re.sub(r'\n+','\n',textbody)  
    return textbody
    
def read_image_ocr(filename):
    '''Retrieves images from PDF and then uses tesseract to extract any visible text from them'''
    tools = pyocr.get_available_tools()
    if len(tools) == 0:
        print('OCR not possible - no readers found')
        return
    tool=tools[0] # retrieve the first tool - should be tesseract
    lang=tool.get_available_languages()[0]
    txt = tool.image_to_string(
        Image.open(filename),
        lang='eng',
        builder=pyocr.builders.TextBuilder()
    )
    word_boxes = tool.image_to_string(
        Image.open(filename),
        lang="eng",
        builder=pyocr.builders.WordBoxBuilder()
    )
    line_and_word_boxes = tool.image_to_string(
        Image.open(filename), lang="eng",
        builder=pyocr.builders.LineBoxBuilder()
    )
    return txt
   
    
    
    

def read_html(filename):
    '''extract all text strings from a html file
    arguments are the filename to read. Returns text in one line'''
    textbody = ''
    try:
        rawtext = str(open(filename, 'rb').read())
        rawtext = re.sub(r'\\x..','', rawtext)
        soup = BeautifulSoup(rawtext, 'html.parser')
        try:        
            rep=soup.script.replace_with(' ')
            while rep:
                rep=soup.script.replace_with(' ')
        except:
            pass
        
        textbody = soup.get_text().replace("'", "").replace('"', '')
    except Exception as e:
        print('Error opening HTML file %s: %s'%(filename, e))
    return textbody

def read_text(filename):
    '''extract text from text file.'''
    return open(filename).read()

def process_file(filename, ppt_images=False, pdf_images=False): #TODO
    '''extracts the keywords from the text for a file.'''

    textbody = ''
    if filename.lower()[-4:] == 'pptx':
        textbody = read_pptx(filename, images=ppt_images)
    elif filename.lower()[-4:] == 'docx':
        textbody = read_docx(filename)
    elif filename.lower()[-3:] == 'pdf':
        textbody = read_pdf(filename, do_ocr=pdf_images)
    elif filename.lower()[-4:] == 'html':
        textbody = read_html(filename)
    elif filename.lower()[-3:] == 'htm':
        textbody = read_html(filename)
    elif filename.lower()[-3:] == 'txt':
        textbody = read_text(filename)
    else:
        print('Incorrect filetype. Cannot parse %s'%filename)
    #keywords = filter(lambda x: (anword.match(x) and len(x) >2),
    #filter(lambda y: y not in stopwords, [x.lower() for x in textbody.split(" ")]))
    return textbody
#    fileroot = filename.split("\\")[-1]
#    eventcode = "_".join([module, week,  fileroot.replace(" ", "_")])
#    sys.stdout.write("%s\t%s\n"%(eventcode, ",".join(sorted(keywords))))

def write_output(outdir, filename, module, week, file, textbody, tfh):
    '''Writes a formatted file to disk containing the annotation for the teaching document
    outdir is the directory in which to save the data
    filename is the filename under which it is to be saved
    module, week and file are the specific document details
    keywords are from the extracted text'''
    fh = open(os.path.join(outdir, filename), 'w')
    fh.write(str(textbody.encode('utf-8')))
    level=module[2]
    semester=module[3]
    tfh.write('\t'.join([filename, level, semester, module, week, file, str(len(textbody))])+'\n')
    fh.close()
    


def read_dir(datasource, datasink, ppt_images=False, pdf_images=False):
    '''takes a source directory (datasource) organsied as module/week
    and parses all valid files underneath, storing the results as annotated
    keyword lists in datasink'''
    try:
        tfh = open(os.path.join(datasink, 'index.dat'),'w')
        modules = os.listdir(datasource)
        for m in modules:
            if os.path.isdir(os.path.join(datasource, m)):
                module = m
                mfilecount = 0
                try:
                    weeks = os.listdir(os.path.join(datasource, m))
                    for w in weeks:
                        if os.path.isdir(os.path.join(datasource, m, w)):
                            try:
                                files = os.listdir(os.path.join(datasource, m, w))
                                for f in files:
                                    if f.lower()[-4:] == 'pptx' \
                                            or f.lower()[-4:] == 'docx' \
                                            or f.lower()[-3:] == 'pdf' \
                                            or f.lower()[-3:] == 'txt' \
                                            or f.lower()[-4:] == 'html' \
                                            or f.lower()[-3:] == 'htm':
                                        filepath = os.path.join(datasource, m, w, f) 
                                        textbody = str(process_file(filepath, ppt_images=ppt_images, pdf_images=pdf_images).encode('utf-8', 'ignore'))
                                        mfilecount += 1
                                        fname = '%s_%s.txt'%(module, mfilecount)
                                        write_output(datasink, fname, m, w, f, textbody,tfh)
                            except Exception as e:
                                print('Error reading file %s in week %s of module %s: %s'%(f, w, m, e))
                except Exception as e:
                    print('Error reading events in module %s: %s'%(m, e))
        tfh.close()
    except Exception as e:
        raise Exception('Bad data directory: %s'%e)
#parser=argparse.ArgumentParser(description='Learning event text miner')
#parser.add_argument('--datadir', '-d', help='Input event file')
#parser.add_argument('--outfile', '-o', help='output event link file')
#args = parser.parse_args()


#include argument parsing.

class NullStemmer(nltk.stem.api.StemmerI):
    def stem(self, word):
        return word
    def stem_word(self, word, i=0, j=None):
        if j:
            return word[i:j+1]
        else:
            return word

class IndexedText(object):
    def __init__(self, stemmer, text):
        if stemmer == None:
            stemmer = NullStemmer()
        self._text = text
        self._stemmer = stemmer
        self._index = nltk.Index((self._stem(word),i) for (i, word) in enumerate(text))
    def concordance(self, word, width=40):
        matches = []
        key = self._stem(word)
        wc = width/4
        for i in self._index[key]:
            left = int(max(0, i - wc))
            right = int(min(i + wc, len(self._text)))
            #print('left: %s right %s' % (left, right))
            lcontext = ' '.join(self._text[left:i])
            rcontext = ' '.join(self._text[i:right])
            ldisplay = '%*s' % (width, lcontext[-width:])
            rdisplay = '%-*s' % (width, rcontext[:width])
            matches.append(' '.join([ldisplay,rdisplay]))
        return matches
        
    def _stem(self, word):
        return self._stemmer.stem(word).lower()

class TeachingCorpus():
    '''Object oriented way to represent the corpus'''
    def __init__(self, datadir, stemmer=None, nostop=True, bigraph=False, case=False, savefile=None):
        self.keywords={} #for IDF for terms
        self.lectures={} #capture TF values
        self.idf={}
        self.indices={}
        self.tfidf={}
        self._stemmer = stemmer
        if stemmer == None:
            self._stemmer = NullStemmer()
        self._datadir = datadir
        self.corpus = None
        self._nostop = nostop
        self._bigraph = bigraph
        self._case = case
        self._savefile = savefile
        self.corpus = self.load_corpus()
        
        
    def get_keywords(self, tokens):
        '''simple keywords splitting'''
        #stem words    
        anword = re.compile(r'^[0-9a-z-_]+[a-z][0-9a-z-_]+$')
        stopwords = STOPWORDS.words("english")
        if not self._case:
            tokens = [t.lower() for t in tokens]
        if self._stemmer != None:
            tokens = [self._stemmer.stem(t) for t in tokens if t]
        
        doubles = []
        if self._bigraph:
            doubles = list(ngrams(tokens,2))
        kw = tokens
        if self._nostop:    
            kw = list(filter(lambda x: (anword.match(x) and len(x) > 2),
                    filter(lambda y: y not in stopwords, tokens)))
            doubles = filter(lambda x: x[0] not in stopwords and anword.match(x[0]) and\
                                       x[1] not in stopwords and anword.match(x[1]), doubles)
        kw = kw + list([' '.join([str(k) for k in list(w)]) for w in doubles] )                             
        return kw
    
    def get_polywords(self, text):
        '''more complex phrase isolation between stopwords'''
        kw = [x.lower() for x in text.split(" ")]
        for i in enumerate(kw):
            if i[1] in stopwords:
                kw[i[0]] = ","
        kp = [x.strip() for x in " ".join(kw).split(",")]
        return kp

    def do_idf(self, termlist):
        '''corpus keyword prevalence'''
        for t in termlist:
            try:
                self.keywords[t]+=1
            except:
                self.keywords[t]=1
                
    def do_tf(self,key, words):
        '''calculate term frequency for a document with key'''
        tf={}
        for t in words:
            try:
                tf[t]+=1
            except:
                tf[t]=1
        self.lectures[key] = {}
        tc=len(words)
        for k in tf.keys():
            self.lectures[key][k] = float(tf[k])/tc
    
    def load_corpus(self):
        '''Load the files in datadir as a Plantext corpus'''
        try:
            corpus = nltk.corpus.PlaintextCorpusReader(self._datadir,'.*\.txt')
            if self._savefile and os.path.exists(self._savefile):
                #load from save file
                #should be a pickle string for each major dictionary
                sfh = open(self._savefile)
                savedata = pickle.loads(base64.b64decode(sfh.read()))
                sfh.close()
                self.fileinfo = savedata['fileinfo']
                self.idf = savedata['idf']
                self.tfidf = savedata['tfidf']
                self.keywords = savedata['keywords']
                self.lectures = savedata['lectures']                         
                print('Indexing data')                
                for f in corpus.fileids():

                    try:
                        self.indices[f] = IndexedText( self._stemmer, corpus.words(f))
                    except Exception as e:
                        print('Error indexing corpus file %s after save: %s'%(f,e))
            else:                
                for f in corpus.fileids():
                    try:
                        kw = self.get_keywords(corpus.words(f))
                        self.do_idf(set(kw))
                        self.do_tf(f, kw)
                        #build indexes for concordancing
                        self.indices[f] = IndexedText( self._stemmer, corpus.words(f))
                    except Exception as ex:
                        print('Error getting keywords from corpus document %s: %s'%(f,ex))
                kwlen=float(len(self.lectures.keys()))
                for k in self.keywords.keys():
                    self.idf[k] = kwlen / self.keywords[k]
                self.load_fileinfo(self._datadir, filename='index.dat')
                self._do_tfidf()
            return corpus
        except Exception as e:
            print('Error reading corpus from %s: %s'%(self._datadir, traceback.print_exc()))
    
    def load_fileinfo(self, datadir, filename='index.dat'):
        '''loads info abotu each filename from a tab separated file in the data directory.
        By default this is index.dat and it has the following structure:
        data filename
        level
        semester
        module
        week
        source filename
        textfile size
        filetype
        '''
        try:
            self.fileinfo={}
            fh = open(os.path.join(datadir,filename))
            for line in fh.readlines():
                f = line.strip().split('\t')
                if len(f) < 5:
                    continue
                ftype = f[5].split('.')[-1].lower()
                self.fileinfo[f[0]] = dict(zip(
                ('datafile','level','semester','module','week','source','filesize','filetype'), 
                f+[ftype]))                
        except Exception as e:
             print('Error reading fileinfo from %s (%s): %s'%(filename,line, e))
             

    def _do_tfidf(self):
        '''Calculate tf-idf terms for each document.'''
        doclist = list(self.lectures.keys())
        #print('documents %s'%doclist)
        for i in range(len(doclist)):
            if len(list(self.lectures[doclist[i]].keys())) > 0:
                self.tfidf[doclist[i]] = {}
                for k in list(self.lectures[doclist[i]].keys()):
                    self.tfidf[doclist[i]][k] = self.lectures[doclist[i]][k] * self.idf[k]
    
    # now do a pairwise cosine similarity between the documents and output.
    def cal_network(self, outfile, minsize=500):
        '''calculate a cosine similarity network using TF/IDF between all documents in the corpus.
        The results are written to outfile as a tab delimited table.'''   
        if outfile.lower()[-4:] != '.csv':
            outfile += '.csv'
        ofh = open(outfile, 'w')
        validdoclist=list([ f for f in self.tfidf.keys() if int(self.fileinfo[f]['filesize']) > minsize]) # docs with >0 keywords
        for i in range(len(validdoclist)):
            for j in range(i+1,len(validdoclist)):
                try:
                    cs = cossim(self.tfidf[validdoclist[i]], 
                                                self.tfidf[validdoclist[j]])
                    ofh.write(",".join([validdoclist[i], validdoclist[j], str(cs*100000)])+"\n")
                except:
                    print('Error determining cosine value for %s and %s'%(validdoclist[i], validdoclist[j]))
        ofh.close()
        
    def concordance(self, word):
        '''search each document for instances of the particular *word* (including 
        similar terms) and return the annotated list. '''
        occurrences={}
        for p in self.indices.keys():
            matches = self.indices[p].concordance(word)
            if len(matches) > 0:
                occurrences[p] = matches 
        return occurrences
        
    def show_matches(self, word, details=True, lines=False):
        '''get matches and display neatly.
        details - display a number of matches if False.
        Display a header and all matches if True'''
        matches=self.concordance(word)
        print('%s documents returned matches' % len(matches))
        bymodule={}
        for p in sorted(matches.keys(), key = lambda x : len(matches[x]), reverse = True):
            try:
                bymodule[self.fileinfo[p]['module']].append(p)
            except:
                bymodule[self.fileinfo[p]['module']] = [p]  
        for m in sorted(bymodule.keys(), key = lambda x : len(bymodule[x]), reverse=True):
            print('%s (%s files match)'%(m, len(bymodule[m])))
            if details:                 
                for p in bymodule[m]:
                    print("%s %s %s (%s matches)"%(self.fileinfo[p]['module'],
                      self.fileinfo[p]['week'], self.fileinfo[p]['source'], 
                      len(matches[p])))
                    if lines:                   
                        for d in matches[p]:
                            print(d)

    def match_text(self, textkw, matches=10, excludepdf=False, excludemodule=None, minsize=500):
        '''Find the best matches to the given file'''
        kwcount={}
        for k in textkw:
            try:
                kwcount[k] += 1
            except:
                kwcount[k] = 1
        thistfidf=dict([(k, kwcount[k]*self.idf[k]/len(kwcount)) for k in kwcount.keys()])        
                
        scores = {}
        validdoclist=list([f for f in self.tfidf.keys() if int(self.fileinfo[f]['filesize'])> minsize]) # docs with >0 keywords
        if excludepdf:
            validdoclist = [f for f in validdoclist if self.fileinfo[f]['filetype'] != 'pdf']
        if excludemodule:
            validdoclist = [f for f in validdoclist if self.fileinfo[f]['module'].upper() != excludemodule.upper() ]
        for i in range(len(validdoclist)):
            scores[validdoclist[i]] =1000000*cossim(self.tfidf[validdoclist[i]], thistfidf) 
        topscores = sorted(scores.keys(), key=lambda x: scores[x], reverse=True)
        return [(a, scores[a]) for a in topscores[:matches]]

    def match_file(self, filename, matches=20, excludepdf=False, minsize=500):
        '''Extract text from the given file and match text, producing a ranked list of matches'''
        text=process_file(filename)
        
        kw = [k for k in self.get_keywords(nltk.word_tokenize(text)) if k in self.idf.keys()]
        hsm = self.match_text(kw, matches, excludepdf, excludemodule=False, minsize=minsize) 
        print('Your file %s matched the following modules/documents:'%filename)
        for hit in hsm:
            fi = self.fileinfo[hit[0]]
            print("%0.d\t%s\t%s\t%s\t%s"%(hit[1],fi['module'],fi['week'], hit[0], fi['source']))
            
    def match_module(self, modulename, matches=5, excludepdf=False, minsize=500):
        '''match results for all files in a module'''
        files=[f for f in self.fileinfo.keys() if self.fileinfo[f]['module'].upper()==modulename.upper() and int(self.fileinfo[f]['filesize']) > minsize]
        for f in sorted(files, key=lambda x: self.fileinfo[x]['week']):
            kw =  self.get_keywords(self.corpus.words(f))
            hsm = self.match_text(kw, matches, excludepdf, excludemodule=modulename, minsize=minsize) 
            print('Your file %s from %s matched the following modules/documents:'%(self.fileinfo[f]['source'],
                                                                                   self.fileinfo[f]['week']))
            for hit in hsm:
                fi = self.fileinfo[hit[0]]
                print("%0.d\t%s\t%s\t%s\t%s"%(hit[1],fi['module'],fi['week'], hit[0], fi['source']))
            print('')
    
    def filematch_info(self, file1, file2):
        '''Performs a TF-IDF match and returns the keywords and pairs for calculation of cosine similarity'''
        kvp = {}
        try:
            kws=set(list(self.tfidf[file1].keys()) + list(self.tfidf[file2].keys()))
            for k in kws:
                t1 = self.tfidf[file1].get(k,0)
                t2 = self.tfidf[file2].get(k,0)
                dp = math.sqrt(t1**2 + t2**2)
                kvp[k] = (t1,t2, dp, (t1**2)/dp, (t2**2)/dp)
        except Exception as e:
            print('Some sort of error occurred matching %s and %s: %s'%(file1, file2, e))
        return kvp
        
    def top_kw(self, file1, file2, tophits=20):
        '''provide a sorted list of keywords with strongest terms in both/one, or the other of file1 and file2'''
        kwlist = self.filematch_info(file1, file2)
        scale = 1000000
        # now sort each way and take the top n scores(scaled appropriately)
        print('Comparison of term TD-IDF scores found in files %s and %s'%(file1, file2))
        print('Top %s terms strongly represented in both files'%tophits)
        for k in sorted(kwlist.keys(), key = lambda x: kwlist[x][2])[:tophits]:
            print('%0.d\t%s'%(kwlist[k][2]*scale, k))
        print('Top %s terms strongly represented in %s but not in %s'%(tophits, file1, file2))
        for k in sorted(kwlist.keys(), key = lambda x: kwlist[x][3])[:tophits]:
            print('%0.d\t%s'%(kwlist[k][3]*scale, k))
        print('Top %s terms strongly represented in %s but not in %s'%(tophits, file2, file1))
        for k in sorted(kwlist.keys(), key = lambda x: kwlist[x][4])[:tophits]:
            print('%0.d\t%s'%(kwlist[k][4]*scale, k))
            
    def save_corpus_to_file(self, filename=None):
        ''' writes the corpus object to file filename'''
        if self._savefile and not filename:
            filename = self._savefile
        elif filename:
            self._savefile = filename
        else:
            print('No filename specified.')
            return
        try:
            # change this to build a file of pickles, don't include the indices.
            dumpstruct ={'fileinfo': self.fileinfo,
                        'idf': self.idf,
                        'tfidf': self.tfidf,
                        'lectures': self.lectures,
                        'keywords': self.keywords}
            #load from save file
            #should be a pickle string for each major dictionary
            sfh = open(filename, 'wb')
            sfh.write(base64.b64encode(pickle.dumps(dumpstruct)))
            sfh.close()
        except Exception as e:
            print('Error saving corpus to file %s: %s'%(filename, e))
        
        

    
    
def dotprod(a, b):
    """ Compute dot product
    Args:
        a (dictionary): first dictionary of record to value
        b (dictionary): second dictionary of record to value
    Returns:
        dotProd: result of the dot product with the two input dictionaries
    """
    prod=0.0
    for k in a.keys():
        try:
            prod += b[k]*a[k]
        except Exception as e:
            pass
    return prod

def norm(a):
    """ Compute square root of the dot product
    Args:
        a (dictionary): a dictionary of record to value
    Returns:
        norm: a dictionary of tokens to its TF values
    """
    sq=0.0
    for k in a.keys():
        sq += a[k]**2
    
    return math.sqrt(sq) 

def cossim(a, b):
    """ Compute cosine similarity
    Args:
        a (dictionary): first dictionary of record to value
        b (dictionary): second dictionary of record to value
    Returns:
        cossim: dot product of two dictionaries divided by the norm of the first dictionary and
                then by the norm of the second dictionary
    """
    try:
        cs = dotprod(a,b)/(norm(a)*norm(b))
        return cs
    except:
        print('len(a) %s len(b) %s dotproduct: %s norm(a) %s norm(b) %s'%(len(a), len(b), dotprod(a,b), norm(a), norm(b)))
        